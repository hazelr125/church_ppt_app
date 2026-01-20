# generate_ppt.py
from pptx import Presentation
from pptx.util import Pt
import re

import fnmatch
FONT_STYLES = {
    "{ANNOUNCEMENTS_TEXT}": {"font_name": "Calibri (MS)", "font_size_pt": 24},
    "{PSALMS_DES}": {"font_name": "Times New Roman MT", "font_size_pt": 60},
    "{OT_DES}": {"font_name": "Times New Roman MT", "font_size_pt": 60},
    "{NT_DES}": {"font_name": "Times New Roman MT", "font_size_pt": 60},
    "{GOSPEL_DES}": {"font_name": "Times New Roman MT", "font_size_pt": 60},
    "PSALM_EN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 50}, 
    "OT_EN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 50}, 
    "NT_EN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 50}, 
    "GOSPEL_EN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 50}, 
    "HYMN*_DES": {"font_name": "Times New Roman MT", "font_size_pt": 60}, 
    #"HYMN*_KN_V*": {"font_name": "Noto Sans Kannada MT", "font_size_pt": 55},
    "HYMN*_KN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 55},
    "HYMN*_EN_V*": {"font_name": "Calibri (MS)", "font_size_pt": 55},           
}

def get_font_style_for_placeholder(ph):
    """Get font style dictionary for a placeholder, supporting wildcards."""
    if ph in FONT_STYLES:
        return FONT_STYLES[ph]
    for pattern, style in FONT_STYLES.items():
        if "*" in pattern and fnmatch.fnmatch(ph.strip("{}"), pattern):
            return style
    return {}

def set_text_frame_text(tf, text, font_name=None, font_size_pt=None):
    tf.clear()
    lines = text.splitlines() if text else ['']
    for i, ln in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ln
        else:
            p = tf.add_paragraph()
            p.text = ln
        for run in p.runs:
            if font_name:
                run.font.name = font_name
            if font_size_pt:
                run.font.size = Pt(font_size_pt)

# Simple check for Kannada characters
KANNADA_UNICODE_RANGE = re.compile(r'[\u0C80-\u0CFF]')

def replace_placeholders(prs, mapping, kannada_font="Noto Sans Kannada", english_font="Arial"):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            txt = shape.text_frame.text
            for ph, val in mapping.items():
                if ph in txt:
                    style = get_font_style_for_placeholder(ph)

                    # Choose font based on style dict, fallback to language
                    font_name = style.get(
                        "font_name",
                        kannada_font if KANNADA_UNICODE_RANGE.search(val or '') else english_font
                    )
                    font_size_pt = style.get("font_size_pt", None)

                    set_text_frame_text(shape.text_frame, val or "", font_name=font_name, font_size_pt=font_size_pt)
                    break  # stop checking more placeholders for this shape

def insert_announcements_table(slide, left, top, width, height, table_data):
    if not table_data:
        return
    rows = len(table_data)
    cols = len(table_data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r, row in enumerate(table_data):
        for c, cell in enumerate(row):
            table.cell(r, c).text = str(cell)

def generate_presentation(template_pptx, output_pptx, mapping, announcements_table_data=None, kannada_font='Noto Sans Kannada'):
    prs = Presentation(template_pptx)
    # 1) basic replacement
    replace_placeholders(prs, mapping, kannada_font)

    # 2) replace ANNOUNCEMENTS_TABLE placeholders by creating a table at same position
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            if '{ANNOUNCEMENTS_TABLE}' in shape.text_frame.text:
                left = shape.left; top = shape.top; width = shape.width; height = shape.height
                shape.text_frame.clear()
                insert_announcements_table(slide, left, top, width, height, announcements_table_data)

    prs.save(output_pptx)